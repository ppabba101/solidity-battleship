"""
_render_preview.py — approximate PNG previews of battleship-zk-demo.pptx.

Since LibreOffice isn't available, we walk the pptx shape tree and rasterize
each slide with Pillow at 1600x900. This is an approximation — font metrics and
some shape effects won't match Keynote/PowerPoint pixel-perfectly — but it's
enough to catch:
  - text overflowing its bounding box
  - shapes colliding / misaligned
  - blank slides or missing elements
  - color / contrast issues

Output: presentation/preview/slide-{n}.png
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image, ImageDraw, ImageFont

ROOT = Path(__file__).parent
PPTX = ROOT / "battleship-zk-demo.pptx"
OUT  = ROOT / "preview"
OUT.mkdir(exist_ok=True)

CANVAS_W, CANVAS_H = 1600, 900  # 16:9

# Fonts: try Inter, fall back to Helvetica, then default
def find_font(names, size):
    for n in names:
        for path in [
            f"/System/Library/Fonts/Supplemental/{n}.ttf",
            f"/System/Library/Fonts/{n}.ttf",
            f"/System/Library/Fonts/{n}.ttc",
            f"/Library/Fonts/{n}.ttf",
            f"/Users/pranavpabba/Library/Fonts/{n}.ttf",
            f"/Users/pranavpabba/Library/Fonts/{n}.otf",
        ]:
            if Path(path).exists():
                try:
                    return ImageFont.truetype(path, size)
                except Exception:
                    pass
    try:
        return ImageFont.truetype("Helvetica", size)
    except Exception:
        return ImageFont.load_default()

BODY_NAMES = ["Inter", "Inter-Regular", "HelveticaNeue", "Helvetica", "Arial"]
BODY_BOLD  = ["Inter-Bold", "HelveticaNeue-Bold", "Helvetica-Bold", "Arial Bold"]
CODE_NAMES = ["JetBrainsMono-Regular", "Menlo", "Monaco", "Courier New"]


SLIDE_W_EMU = 13.333 * 914400  # 12192000
SLIDE_H_EMU = 7.5 * 914400     # 6858000

def emu_to_px_x(v):
    return int(v / SLIDE_W_EMU * CANVAS_W)

def emu_to_px_y(v):
    return int(v / SLIDE_H_EMU * CANVAS_H)

def emu_font_size(pt):
    # Pillow uses pixel sizes. 1pt ~ 1.33 px on a 96dpi screen; we're
    # rendering at 1600/13.333in ≈ 120 px/in ≈ 1.66 px/pt.
    return max(6, int(round(pt * 1.66)))


def rgb_from_run(run):
    try:
        c = run.font.color.rgb
        if c is None:
            return (248, 250, 252)
        return (c[0], c[1], c[2])
    except Exception:
        return (248, 250, 252)


def shape_fill(shape):
    try:
        f = shape.fill
        if f.type == 1:  # solid
            c = f.fore_color.rgb
            return (c[0], c[1], c[2])
    except Exception:
        pass
    return None


def shape_line(shape):
    try:
        c = shape.line.color.rgb
        if c is None:
            return None
        return (c[0], c[1], c[2])
    except Exception:
        return None


def draw_slide(slide, idx):
    img = Image.new("RGB", (CANVAS_W, CANVAS_H), (10, 15, 26))
    draw = ImageDraw.Draw(img)

    overflow_warnings = []

    for shape in slide.shapes:
        try:
            L = emu_to_px_x(shape.left or 0)
            T = emu_to_px_y(shape.top or 0)
            W = emu_to_px_x(shape.width or 0)
            H = emu_to_px_y(shape.height or 0)
        except Exception:
            continue

        # Geometric shapes (rectangles, ovals, rounded rects, lines)
        if shape.shape_type in (MSO_SHAPE_TYPE.AUTO_SHAPE,):
            fill = shape_fill(shape)
            outline = shape_line(shape)
            # Detect oval vs rectangle vs rounded-rectangle
            try:
                geom = str(shape.auto_shape_type)
            except Exception:
                geom = ""
            if "OVAL" in geom:
                if fill:
                    draw.ellipse([L, T, L + W, T + H], fill=fill, outline=outline)
            elif "ROUNDED" in geom:
                r = min(W, H) // 6
                if fill:
                    draw.rounded_rectangle([L, T, L + W, T + H], radius=r,
                                           fill=fill, outline=outline)
                elif outline:
                    draw.rounded_rectangle([L, T, L + W, T + H], radius=r,
                                           outline=outline)
            else:
                if fill:
                    draw.rectangle([L, T, L + W, T + H], fill=fill, outline=outline)
                elif outline:
                    draw.rectangle([L, T, L + W, T + H], outline=outline)

        # Text boxes
        if shape.has_text_frame:
            tf = shape.text_frame
            y = T
            for p in tf.paragraphs:
                # Combine runs in a paragraph
                line_text = ""
                max_size = 12
                color = (248, 250, 252)
                bold = False
                mono = False
                for run in p.runs:
                    line_text += run.text or ""
                    try:
                        if run.font.size is not None:
                            sz = int(run.font.size.pt)
                            if sz > max_size:
                                max_size = sz
                    except Exception:
                        pass
                    try:
                        if run.font.bold:
                            bold = True
                    except Exception:
                        pass
                    try:
                        if (run.font.name or "").lower().startswith(("jet", "mon", "con", "cou")):
                            mono = True
                    except Exception:
                        pass
                    c = rgb_from_run(run)
                    if c != (248, 250, 252):
                        color = c
                if not line_text.strip():
                    y += emu_font_size(max_size) + 4
                    continue
                names = (CODE_NAMES if mono
                         else (BODY_BOLD if bold else BODY_NAMES))
                font = find_font(names, emu_font_size(max_size))
                # Alignment
                align = p.alignment
                try:
                    bbox = draw.textbbox((0, 0), line_text, font=font)
                    tw = bbox[2] - bbox[0]
                    th = bbox[3] - bbox[1]
                except Exception:
                    tw = len(line_text) * max_size // 2
                    th = max_size
                if align and "CENTER" in str(align):
                    tx = L + (W - tw) // 2
                elif align and "RIGHT" in str(align):
                    tx = L + W - tw
                else:
                    tx = L
                draw.text((tx, y), line_text, fill=color, font=font)
                # Overflow detection
                if tw > W + 8:
                    overflow_warnings.append(
                        f"slide {idx}: '{line_text[:40]}' text_w={tw}px box_w={W}px"
                    )
                y += th + 6

    # Save
    out = OUT / f"slide-{idx:02d}.png"
    img.save(out)
    return out, overflow_warnings


def main():
    prs = Presentation(str(PPTX))
    all_warns = []
    for i, slide in enumerate(prs.slides, 1):
        path, warns = draw_slide(slide, i)
        print(f"rendered {path.name}")
        all_warns.extend(warns)
    if all_warns:
        print("\nOVERFLOW WARNINGS:")
        for w in all_warns:
            print(" -", w)
    else:
        print("\nNo text overflow detected.")


if __name__ == "__main__":
    main()
