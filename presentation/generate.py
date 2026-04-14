"""
generate.py — builds battleship-zk-demo.pptx via python-pptx.

Usage:
    python3 generate.py

Output:
    battleship-zk-demo.pptx  (same directory as this script)

Design system
-------------
Background:   #0A0F1A  (near-black navy)
Surface:      #1B2942  (elevated panels)
Hairline:     #273449
Primary text: #F8FAFC
Secondary:    #94A3B8
Accent:       #F97316  (orange)
Secondary accent: #38BDF8 (cyan)
Warning:      #F87171  (soft red)

Typography:   Inter (body) / JetBrains Mono (code), with OS fallbacks

Glyph policy
------------
We stick to ASCII and a small whitelist of unicode symbols with broad font
coverage across Inter / Helvetica Neue / Calibri / system defaults:

    ->          U+2192  right arrow
    *           bullet  (we use "-" instead)
    YES / NO    plain words instead of check/ballot
    U+2714 / U+2718  check / ballot — used sparingly, only where
                     surrounding text is clearly labelled so a missing
                     glyph still reads.

The play triangle U+25B6 is REPLACED with the text "LIVE DEMO ->" to
avoid glyph fallback boxes on default installs.
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ---------------------------------------------------------------------------
# Palette
# ---------------------------------------------------------------------------
BG          = RGBColor(0x0A, 0x0F, 0x1A)
SURFACE     = RGBColor(0x1B, 0x29, 0x42)
SURFACE_HI  = RGBColor(0x22, 0x33, 0x52)
HAIRLINE    = RGBColor(0x27, 0x34, 0x49)
TEXT        = RGBColor(0xF8, 0xFA, 0xFC)
MUTED       = RGBColor(0x94, 0xA3, 0xB8)
ORANGE      = RGBColor(0xF9, 0x73, 0x16)
CYAN        = RGBColor(0x38, 0xBD, 0xF8)
WARN        = RGBColor(0xF8, 0x71, 0x71)
GOOD        = RGBColor(0x4A, 0xDE, 0x80)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)

# ---------------------------------------------------------------------------
# Typography
# ---------------------------------------------------------------------------
BODY_FONT = "Inter"
CODE_FONT = "JetBrains Mono"

# ---------------------------------------------------------------------------
# Canvas
# ---------------------------------------------------------------------------
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN  = Inches(0.6)
TOTAL_SLIDES = 11

# Safe unicode we actually use
ARROW = "\u2192"   # ->  right arrow, nearly universal
MDASH = "\u2014"   # em dash
BULL  = "\u2022"   # bullet (middle dot is U+00B7 which also works)
MIDOT = "\u00b7"   # middle dot

# ---------------------------------------------------------------------------
# Primitives
# ---------------------------------------------------------------------------

def new_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs: Presentation, bg=BG):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bgrect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bgrect.line.fill.background()
    bgrect.fill.solid()
    bgrect.fill.fore_color.rgb = bg
    return slide


def add_rect(slide, left, top, width, height,
             fill_color=None, line_color=None, line_width=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    if fill_color is not None:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color is not None:
        shape.line.color.rgb = line_color
        if line_width is not None:
            shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape


def add_round_rect(slide, left, top, width, height,
                   fill_color=None, line_color=None, line_width=None,
                   corner=0.12):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    try:
        shape.adjustments[0] = corner
    except Exception:
        pass
    if fill_color is not None:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color is not None:
        shape.line.color.rgb = line_color
        if line_width is not None:
            shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape


def add_oval(slide, left, top, width, height, fill_color=None, line_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, width, height)
    if fill_color is not None:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color is not None:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, left, top, width, height, text,
             font=BODY_FONT, size=18, bold=False, color=TEXT,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             tracking=0):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    if tracking:
        from pptx.oxml.ns import qn
        rPr = run._r.get_or_add_rPr()
        rPr.set("spc", str(tracking))
    return tb, tf


# ---------------------------------------------------------------------------
# Chrome
# ---------------------------------------------------------------------------

def add_wordmark(slide):
    tb = slide.shapes.add_textbox(MARGIN, Inches(7.0), Inches(3.0), Inches(0.3))
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    r1 = p.add_run()
    r1.text = "battleship"
    r1.font.name = CODE_FONT
    r1.font.size = Pt(10)
    r1.font.color.rgb = MUTED
    r1.font.bold = False
    r2 = p.add_run()
    r2.text = ".zk"
    r2.font.name = CODE_FONT
    r2.font.size = Pt(10)
    r2.font.color.rgb = ORANGE
    r2.font.bold = True


def add_progress_bar(slide, n, total=TOTAL_SLIDES):
    bar_top = Inches(7.35)
    bar_h   = Pt(3)
    bar_l   = MARGIN
    bar_w   = SLIDE_W - MARGIN * 2
    add_rect(slide, bar_l, bar_top, bar_w, bar_h, fill_color=HAIRLINE)
    frac = n / total
    add_rect(slide, bar_l, bar_top, Emu(int(bar_w * frac)), bar_h, fill_color=ORANGE)
    add_text(
        slide,
        SLIDE_W - MARGIN - Inches(1.2), Inches(7.0),
        Inches(1.2), Inches(0.3),
        f"{n:02d} / {total:02d}",
        font=CODE_FONT, size=10, color=MUTED, align=PP_ALIGN.RIGHT,
    )


def add_eyebrow(slide, text, top=Inches(0.9)):
    add_text(
        slide,
        MARGIN, top, Inches(12.0), Inches(0.35),
        text.upper(),
        font=BODY_FONT, size=11, bold=True, color=CYAN,
        tracking=300,
    )


def add_title(slide, text, top=Inches(1.25)):
    add_text(
        slide,
        MARGIN, top, Inches(12.0), Inches(0.9),
        text,
        font=BODY_FONT, size=36, bold=True, color=TEXT,
    )


def add_rule(slide, top=Inches(2.1), width=Inches(1.1)):
    add_rect(slide, MARGIN, top, width, Pt(3), fill_color=ORANGE)


def chrome(slide, n, eyebrow=None, title=None):
    if eyebrow:
        add_eyebrow(slide, eyebrow)
    if title:
        add_title(slide, title)
        add_rule(slide)
    add_wordmark(slide)
    add_progress_bar(slide, n)


# ---------------------------------------------------------------------------
# Flow helpers
# ---------------------------------------------------------------------------

def add_flow_station(slide, left, top, width, height,
                     glyph, title, subtitle,
                     fill=SURFACE, border=HAIRLINE, glyph_color=ORANGE):
    add_round_rect(slide, left, top, width, height,
                   fill_color=fill, line_color=border, line_width=Pt(1),
                   corner=0.22)
    add_text(
        slide, left + Inches(0.15), top + Inches(0.1),
        Inches(1.2), Inches(0.35),
        glyph,
        font=CODE_FONT, size=11, bold=True, color=glyph_color,
        tracking=150,
    )
    add_text(
        slide, left + Inches(0.15), top + Inches(0.45),
        width - Inches(0.3), Inches(0.45),
        title,
        font=BODY_FONT, size=15, bold=True, color=TEXT,
    )
    add_text(
        slide, left + Inches(0.15), top + Inches(0.92),
        width - Inches(0.3), Inches(0.35),
        subtitle,
        font=CODE_FONT, size=9, bold=False, color=MUTED,
    )


def add_arrow(slide, x1, y1, x2, y2, color=ORANGE):
    conn = slide.shapes.add_connector(2, x1, y1, x2, y2)
    conn.line.color.rgb = color
    conn.line.width = Pt(1.75)
    line_elem = conn.line._get_or_add_ln()
    from pptx.oxml.ns import qn
    from lxml import etree
    tailEnd = etree.SubElement(line_elem, qn("a:tailEnd"))
    tailEnd.set("type", "triangle")
    tailEnd.set("w", "med")
    tailEnd.set("h", "med")
    return conn


def set_notes(slide, notes_text: str):
    tf = slide.notes_slide.notes_text_frame
    tf.text = notes_text


def _draw_bullets(slide, bullets, left, top, width,
                  mark_color=ORANGE, text_color=TEXT):
    row_h = Inches(0.6)
    for i, (mark, text) in enumerate(bullets):
        y = top + row_h * i
        add_text(
            slide, left, y, Inches(0.55), Inches(0.45),
            mark,
            font=BODY_FONT, size=16, bold=True, color=mark_color,
        )
        add_text(
            slide, left + Inches(0.6), y, width - Inches(0.6), Inches(0.5),
            text,
            font=BODY_FONT, size=16, bold=False, color=text_color,
        )


# ===========================================================================
# Slide 1 — Title
# ===========================================================================

def slide_1_title(prs):
    slide = blank_slide(prs)

    add_rect(slide, 0, 0, Inches(0.18), SLIDE_H, fill_color=ORANGE)

    add_text(
        slide,
        Inches(0.9), Inches(1.35),
        Inches(12.0), Inches(0.35),
        f"A ZERO-KNOWLEDGE DEMO   {MIDOT}   BLOCKCHAIN CLUB",
        font=BODY_FONT, size=12, bold=True, color=CYAN, tracking=350,
    )

    add_text(
        slide,
        Inches(0.85), Inches(1.95),
        Inches(12.0), Inches(1.8),
        "Battleship,",
        font=BODY_FONT, size=72, bold=True, color=TEXT,
    )
    add_text(
        slide,
        Inches(0.85), Inches(3.25),
        Inches(12.0), Inches(1.8),
        "proven.",
        font=BODY_FONT, size=72, bold=True, color=ORANGE,
    )

    add_text(
        slide,
        Inches(0.9), Inches(4.85),
        Inches(12.0), Inches(0.5),
        "Eight minutes on how a single zk-SNARK turns a game of hidden state into a trustless one.",
        font=BODY_FONT, size=18, bold=False, color=MUTED,
    )

    # Presenter chip — real names
    chip_w = Inches(5.6)
    chip_h = Inches(0.55)
    chip_l = Inches(0.85)
    chip_t = Inches(6.0)
    add_round_rect(slide, chip_l, chip_t, chip_w, chip_h,
                   fill_color=SURFACE, line_color=HAIRLINE,
                   line_width=Pt(0.75), corner=0.5)
    add_oval(slide, chip_l + Inches(0.22), chip_t + Inches(0.2),
             Inches(0.15), Inches(0.15), fill_color=ORANGE)
    add_text(
        slide,
        chip_l + Inches(0.48), chip_t + Inches(0.12),
        chip_w - Inches(0.6), chip_h,
        f"VIKRAM AKKALA    {MIDOT}    PRANAV PABBA",
        font=BODY_FONT, size=11, bold=True, color=TEXT, tracking=200,
    )

    add_text(
        slide,
        SLIDE_W - MARGIN - Inches(4.5), Inches(6.12),
        Inches(4.5), Inches(0.3),
        f"noir {ARROW} ultraplonk {ARROW} solidity",
        font=CODE_FONT, size=11, bold=False, color=MUTED,
        align=PP_ALIGN.RIGHT,
    )

    add_progress_bar(slide, 1)

    notes = """\
[Pranav]
Welcome everyone. I'm Pranav Pabba and this is my teammate Vikram Akkala. Over the next eight \
minutes we're going to show you Battleship -- the classic two-player board game -- reimagined \
with zero-knowledge proofs. I'll drive the demo at the end; Vikram is going to carry you \
through the cryptography in the middle. By the time we're done you'll understand exactly why \
a naive commitment scheme fails, how Merkle trees get you most of the way, and why a single \
zk-SNARK is the piece that actually closes the gap.

[Vikram]
Thanks Pranav. My job is the crypto -- the problem framing, the Merkle deep-dive, and the \
zk-SNARK reveal. I promise to keep it concrete and tied to Battleship the whole way. Let's go.
"""
    set_notes(slide, notes)
    return slide


# ===========================================================================
# Slide 2 — The Problem
# ===========================================================================

def slide_2_problem(prs):
    slide = blank_slide(prs)
    chrome(slide, 2, eyebrow=f"01  {MIDOT}  The Problem",
           title="A game built on hidden state.")

    add_text(
        slide,
        MARGIN, Inches(2.6),
        Inches(1.2), Inches(2.0),
        "\u201C",
        font=BODY_FONT, size=180, bold=True, color=ORANGE,
    )

    tb = slide.shapes.add_textbox(Inches(1.9), Inches(3.05),
                                  Inches(10.9), Inches(2.6))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, line in enumerate([
        "How do you prove your board",
        f"is legal{MDASH}without revealing it?",
    ]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_before = Pt(0)
        r = p.add_run()
        r.text = line
        r.font.name = BODY_FONT
        r.font.size = Pt(42)
        r.font.bold = True
        r.font.color.rgb = TEXT if i == 0 else ORANGE

    add_text(
        slide,
        Inches(1.9), Inches(5.55),
        Inches(10.9), Inches(0.4),
        f"{MDASH} the central cryptographic puzzle",
        font=BODY_FONT, size=14, bold=False, color=MUTED,
    )

    notes = """\
[Vikram]
Battleship is a game of hidden information. I place ships on my private board, you place ships \
on yours, and neither of us sees the other's grid. That's fine when we're sitting across a \
table with cardboard pegs -- but in a digital implementation, nothing stops me from just \
lying. I can claim every one of your shots is a miss and coast to victory. The rules don't \
enforce themselves. So we need a way for me to commit to a board up front, prove that board is \
a legal fleet, and then answer each of your shots honestly -- all without ever showing you \
where my ships actually are. That is the central cryptographic puzzle of this talk.

[Pranav]
Keep this puzzle in your head -- every slide from here on is chipping away at it.
"""
    set_notes(slide, notes)
    return slide


# ===========================================================================
# Slide 3 — Merkle 101
# ===========================================================================

def slide_3_merkle_101(prs):
    slide = blank_slide(prs)
    chrome(slide, 3, eyebrow=f"02  {MIDOT}  Merkle 101",
           title="A Merkle tree, in one picture.")

    # Draw a small 4-leaf tree on the left
    tree_l = MARGIN
    tree_t = Inches(2.6)
    tree_w = Inches(6.8)
    tree_h = Inches(3.9)

    add_round_rect(slide, tree_l, tree_t, tree_w, tree_h,
                   fill_color=SURFACE, line_color=HAIRLINE,
                   line_width=Pt(1), corner=0.06)

    # Root
    root_cx = tree_l + tree_w / 2
    root_cy = tree_t + Inches(0.55)
    node_w  = Inches(1.6)
    node_h  = Inches(0.45)
    add_round_rect(slide, root_cx - node_w / 2, root_cy,
                   node_w, node_h,
                   fill_color=ORANGE, line_color=ORANGE, corner=0.4)
    add_text(slide, root_cx - node_w / 2, root_cy + Inches(0.05),
             node_w, node_h, "ROOT",
             font=BODY_FONT, size=13, bold=True, color=BG,
             align=PP_ALIGN.CENTER)

    # Level 1 — two nodes
    l1_y = tree_t + Inches(1.4)
    l1_cx_left  = tree_l + tree_w * 0.3
    l1_cx_right = tree_l + tree_w * 0.7
    for cx, label in [(l1_cx_left, "H(AB)"), (l1_cx_right, "H(CD)")]:
        add_round_rect(slide, cx - node_w / 2, l1_y,
                       node_w, node_h,
                       fill_color=SURFACE_HI, line_color=HAIRLINE, corner=0.4)
        add_text(slide, cx - node_w / 2, l1_y + Inches(0.05),
                 node_w, node_h, label,
                 font=CODE_FONT, size=12, bold=True, color=TEXT,
                 align=PP_ALIGN.CENTER)

    # Level 2 — 4 leaves
    l2_y = tree_t + Inches(2.35)
    leaf_cxs = [tree_l + tree_w * f for f in (0.15, 0.38, 0.62, 0.85)]
    leaf_labels = ["cell A", "cell B", "cell C", "cell D"]
    for cx, label in zip(leaf_cxs, leaf_labels):
        add_round_rect(slide, cx - Inches(0.65), l2_y,
                       Inches(1.3), Inches(0.45),
                       fill_color=BG, line_color=CYAN, corner=0.4)
        add_text(slide, cx - Inches(0.65), l2_y + Inches(0.05),
                 Inches(1.3), Inches(0.45), label,
                 font=CODE_FONT, size=11, bold=True, color=CYAN,
                 align=PP_ALIGN.CENTER)

    # Connector lines (approximate — use thin rectangles as stems)
    def stem(x1, y1, x2, y2):
        add_arrow(slide, x1, y1, x2, y2, color=HAIRLINE)

    stem(int(l1_cx_left),  int(l1_y), int(root_cx), int(root_cy + node_h))
    stem(int(l1_cx_right), int(l1_y), int(root_cx), int(root_cy + node_h))
    stem(int(leaf_cxs[0]), int(l2_y), int(l1_cx_left),  int(l1_y + node_h))
    stem(int(leaf_cxs[1]), int(l2_y), int(l1_cx_left),  int(l1_y + node_h))
    stem(int(leaf_cxs[2]), int(l2_y), int(l1_cx_right), int(l1_y + node_h))
    stem(int(leaf_cxs[3]), int(l2_y), int(l1_cx_right), int(l1_y + node_h))

    # Caption under the tree
    add_text(slide, tree_l + Inches(0.3), tree_t + Inches(3.1),
             tree_w - Inches(0.6), Inches(0.6),
             "Hash pairs up the tree. The root fingerprints every leaf.",
             font=BODY_FONT, size=13, bold=False, color=MUTED,
             align=PP_ALIGN.CENTER)

    # Right column — the core property
    right_l = tree_l + tree_w + Inches(0.4)
    right_w = SLIDE_W - right_l - MARGIN

    add_text(slide, right_l, tree_t + Inches(0.1),
             right_w, Inches(0.4),
             "THE CORE PROPERTY",
             font=BODY_FONT, size=11, bold=True, color=CYAN, tracking=250)

    add_text(slide, right_l, tree_t + Inches(0.5),
             right_w, Inches(0.7),
             "One root commits to every leaf.",
             font=BODY_FONT, size=22, bold=True, color=TEXT)

    props = [
        (f"{ARROW}", "Publish only the root on-chain."),
        (f"{ARROW}", "Reveal any one leaf with a short path."),
        (f"{ARROW}", "The root cannot be forged later."),
        (f"{ARROW}", "100 cells means ~7 hashes per proof."),
    ]
    _draw_bullets(slide, props, right_l, tree_t + Inches(1.4),
                  right_w, mark_color=ORANGE, text_color=TEXT)

    notes = """\
[Vikram]
Before we get to the zk-SNARK I want to make sure everyone has Merkle trees in their head, \
because this is the piece of the talk that actually ties the cryptography to Battleship. A \
Merkle tree is what you get when you take a list of values, hash them in pairs, then hash those \
hashes in pairs, and keep going until you collapse the whole thing into a single hash at the \
top. That top hash is the root. The magic property is this: the root is a fingerprint of every \
single leaf at once. If I change any leaf, even one bit, the root changes. So if I publish just \
the root, I have committed to the entire list without revealing any of it. Later, I can prove \
"leaf number 37 was X" by handing you X plus about log-base-two-of-a-hundred -- so seven -- \
sibling hashes along the path. You hash your way back up to the root and check it matches. \
That is it. Every blockchain you have ever heard of uses this trick.
"""
    set_notes(slide, notes)
    return slide


# ===========================================================================
# Slide 4 — Merkle in Battleship
# ===========================================================================

def slide_4_merkle_battleship(prs):
    slide = blank_slide(prs)
    chrome(slide, 4, eyebrow=f"03  {MIDOT}  Merkle in Battleship",
           title="Every cell is a leaf.")

    # Left: a 10x10 mini grid with a couple of shaded cells
    grid_l = MARGIN
    grid_t = Inches(2.55)
    grid_side = Inches(3.9)
    cell = grid_side / 10

    add_round_rect(slide, grid_l - Inches(0.15), grid_t - Inches(0.15),
                   grid_side + Inches(0.3), grid_side + Inches(0.3),
                   fill_color=SURFACE, line_color=HAIRLINE,
                   line_width=Pt(1), corner=0.05)

    ship_cells = {(1, 2), (2, 2), (3, 2), (5, 6), (5, 7)}
    for r in range(10):
        for c in range(10):
            x = grid_l + cell * c
            y = grid_t + cell * r
            fill = ORANGE if (c, r) in ship_cells else BG
            add_rect(slide, x, y, cell, cell,
                     fill_color=fill, line_color=HAIRLINE,
                     line_width=Pt(0.5))

    add_text(slide, grid_l, grid_t + grid_side + Inches(0.15),
             grid_side, Inches(0.4),
             "100 cells = 100 leaves",
             font=BODY_FONT, size=12, bold=True, color=MUTED,
             align=PP_ALIGN.CENTER)

    # Right column — leaf + commit + shot-response steps
    right_l = grid_l + grid_side + Inches(0.7)
    right_w = SLIDE_W - right_l - MARGIN

    add_text(slide, right_l, grid_t - Inches(0.05),
             right_w, Inches(0.4),
             "THE RECIPE",
             font=BODY_FONT, size=11, bold=True, color=CYAN, tracking=250)

    # Step blocks
    steps = [
        ("LEAF",
         "leaf_i = hash(cellIndex, occupied, salt_i)"),
        ("COMMIT",
         "publish merkle_root(leaves[0..100]) on-chain"),
        ("SHOT",
         "opponent fires at (x, y)"),
        ("RESPOND",
         "reveal leaf and 7-hash path under root"),
        ("VERIFY",
         "contract recomputes root -> accepts hit/miss"),
    ]
    sy = grid_t + Inches(0.45)
    for i, (label, body) in enumerate(steps):
        y = sy + Inches(0.55) * i
        add_round_rect(slide, right_l, y, Inches(1.15), Inches(0.38),
                       fill_color=SURFACE, line_color=HAIRLINE, corner=0.3)
        add_text(slide, right_l, y + Inches(0.05),
                 Inches(1.15), Inches(0.35),
                 label,
                 font=CODE_FONT, size=10, bold=True, color=ORANGE,
                 align=PP_ALIGN.CENTER, tracking=150)
        add_text(slide, right_l + Inches(1.3), y + Inches(0.04),
                 right_w - Inches(1.3), Inches(0.4),
                 body,
                 font=CODE_FONT, size=12, bold=False, color=TEXT)

    notes = """\
[Vikram]
Now let us plug this into Battleship. Each of my 100 cells becomes a leaf. The leaf is a hash \
of three things: the cell's index zero through ninety-nine, a bit saying whether it's occupied, \
and a random salt so you can't brute-force guess it. I build the Merkle tree over those 100 \
leaves, and I send you the root. That is my commitment. Now you fire a shot at coordinate \
three-comma-five. I look up leaf number thirty-five, tell you the occupied bit, and hand you \
the seven sibling hashes that form the Merkle path. You hash it back up to my root and check. \
If the root matches, you know I'm not making up the answer -- that cell really was whatever I \
said it was when I committed. And crucially, I never had to show you the other 99 cells. This \
is the exact scheme the original Battleship contract in this repo used. It works. It's pretty. \
It feels like we're done. But we are not done.
"""
    set_notes(slide, notes)
    return slide


# ===========================================================================
# Slide 5 — The Merkle Hole (pivot)
# ===========================================================================

def slide_5_merkle_hole(prs):
    slide = blank_slide(prs)
    chrome(slide, 5, eyebrow=f"04  {MIDOT}  The Merkle Hole",
           title="What Merkle cannot prove.")

    col_top = Inches(2.55)
    col_h   = Inches(3.4)
    gap     = Inches(0.4)
    col_w   = (SLIDE_W - MARGIN * 2 - gap) / 2

    # Left: what Merkle proves
    left_l = MARGIN
    add_round_rect(slide, left_l, col_top, col_w, col_h,
                   fill_color=SURFACE, line_color=HAIRLINE,
                   line_width=Pt(1), corner=0.08)
    add_text(slide, left_l + Inches(0.4), col_top + Inches(0.3),
             col_w - Inches(0.8), Inches(0.4),
             "MERKLE PROVES",
             font=BODY_FONT, size=11, bold=True, color=CYAN, tracking=250)
    add_text(slide, left_l + Inches(0.4), col_top + Inches(0.7),
             col_w - Inches(0.8), Inches(0.55),
             "Consistency.",
             font=BODY_FONT, size=26, bold=True, color=TEXT)
    left_bullets = [
        ("YES", "You committed to 100 values."),
        ("YES", "You cannot change them later."),
        ("YES", "Any single reveal is honest."),
    ]
    _draw_bullets(slide, left_bullets, left_l + Inches(0.4),
                  col_top + Inches(1.55), col_w - Inches(0.8),
                  mark_color=GOOD, text_color=TEXT)

    # Right: what Merkle does NOT prove
    right_l = left_l + col_w + gap
    add_round_rect(slide, right_l, col_top, col_w, col_h,
                   fill_color=RGBColor(0x2A, 0x18, 0x20),
                   line_color=RGBColor(0x5A, 0x29, 0x33),
                   line_width=Pt(1), corner=0.08)
    add_text(slide, right_l + Inches(0.4), col_top + Inches(0.3),
             col_w - Inches(0.8), Inches(0.4),
             "MERKLE CANNOT PROVE",
             font=BODY_FONT, size=11, bold=True, color=WARN, tracking=250)
    add_text(slide, right_l + Inches(0.4), col_top + Inches(0.7),
             col_w - Inches(0.8), Inches(0.55),
             "Legality.",
             font=BODY_FONT, size=26, bold=True, color=WARN)
    right_bullets = [
        ("NO", "That the 100 values form a real fleet."),
        ("NO", "No overlaps, no diagonals, correct shapes."),
        ("NO", "Deferred reveal arrives too late."),
    ]
    _draw_bullets(slide, right_bullets, right_l + Inches(0.4),
                  col_top + Inches(1.55), col_w - Inches(0.8),
                  mark_color=WARN, text_color=TEXT)

    # The attack callout — full width, bottom
    callout_t = Inches(6.2)
    add_round_rect(slide, MARGIN, callout_t, SLIDE_W - MARGIN * 2,
                   Inches(0.85),
                   fill_color=RGBColor(0x33, 0x1A, 0x0B),
                   line_color=ORANGE, line_width=Pt(1.25), corner=0.25)
    add_text(slide, MARGIN + Inches(0.3), callout_t + Inches(0.08),
             SLIDE_W - MARGIN * 2 - Inches(0.6), Inches(0.4),
             f"THE ATTACK   {ARROW}",
             font=BODY_FONT, size=11, bold=True, color=ORANGE, tracking=250)
    add_text(slide, MARGIN + Inches(0.3), callout_t + Inches(0.38),
             SLIDE_W - MARGIN * 2 - Inches(0.6), Inches(0.45),
             "Commit an all-empty board. Answer MISS to every shot. Win.",
             font=CODE_FONT, size=14, bold=True, color=TEXT)

    notes = """\
[Vikram]
Here is the hole, and this is the pivot of the whole talk. A Merkle proof proves that I did not \
change my mind about a cell after I committed. It proves consistency. What it absolutely does \
not prove is that the thing I committed to in the first place was a legal Battleship fleet. \
Nothing in a Merkle tree says "this has exactly one carrier, one battleship, two cruisers, one \
destroyer, no overlaps, no diagonals." And now the exploit writes itself. I commit to an \
all-empty ten-by-ten grid. That is a perfectly valid Merkle tree with a perfectly valid root. \
You start firing at me and I honestly, Merkle-verifiably, answer MISS to every single one of \
your shots -- because every cell really is empty. You never sink a ship, I never take damage, \
and on the last turn I claim the win. The deferred "reveal at the end" check that's supposed to \
catch this is theatre -- by the time we'd run it the game is already lost, and in the original \
contract in this repo both branches of that check actually assign the same winner. So Merkle \
alone is not enough. We need legality proven up front, before a single shot is fired. That is \
where zk-SNARKs come in.
"""
    set_notes(slide, notes)
    return slide


# ===========================================================================
# Slide 6 — Enter zk-SNARKs
# ===========================================================================

def slide_6_zksnarks(prs):
    slide = blank_slide(prs)
    chrome(slide, 6, eyebrow=f"05  {MIDOT}  The primitive",
           title="Enter zk-SNARKs.")

    # Big idea line up top
    add_text(slide, MARGIN, Inches(2.5),
             SLIDE_W - MARGIN * 2, Inches(0.6),
             "The proof IS the validity certificate.",
             font=BODY_FONT, size=26, bold=True, color=ORANGE)

    add_text(slide, MARGIN, Inches(3.05),
             SLIDE_W - MARGIN * 2, Inches(0.5),
             "One proof, generated once at commit time, verified once on-chain.",
             font=BODY_FONT, size=15, bold=False, color=MUTED)

    # 2x2 tile grid — ASCII / english labels instead of math glyphs
    grid_top = Inches(3.75)
    grid_h   = Inches(2.9)
    gap      = Inches(0.3)
    tile_w   = (SLIDE_W - MARGIN * 2 - gap) / 2
    tile_h   = (grid_h - gap) / 2

    tiles = [
        ("PROOF",
         "Expressive predicates",
         "Prove the full fleet constraint in one shot."),
        ("CIRCUIT",
         "board_validity.nr",
         "3,482 constraints: shapes, counts, no diagonals."),
        ("HASH",
         "Pedersen binding",
         "Commitment and legality locked into one proof."),
        ("VERIFY",
         "On-chain, ~250k gas",
         "HonkVerifier.sol auto-generated by Noir."),
    ]

    for i, (glyph, title, body) in enumerate(tiles):
        row, col = divmod(i, 2)
        x = MARGIN + (tile_w + gap) * col
        y = grid_top + (tile_h + gap) * row
        add_round_rect(slide, x, y, tile_w, tile_h,
                       fill_color=SURFACE, line_color=HAIRLINE,
                       line_width=Pt(1), corner=0.1)
        add_text(slide, x + Inches(0.3), y + Inches(0.22),
                 Inches(1.9), Inches(0.45),
                 glyph,
                 font=CODE_FONT, size=13, bold=True, color=ORANGE, tracking=200)
        add_text(slide, x + Inches(0.3), y + Inches(0.55),
                 tile_w - Inches(0.6), Inches(0.5),
                 title,
                 font=BODY_FONT, size=18, bold=True, color=TEXT)
        add_text(slide, x + Inches(0.3), y + Inches(0.9),
                 tile_w - Inches(0.6), Inches(0.6),
                 body,
                 font=BODY_FONT, size=13, bold=False, color=MUTED)

    notes = """\
[Vikram]
Here is the move. Instead of proving a single cell's consistency at shot time, we prove an \
entire CONSTRAINT at commit time: "the 100 private cells I am about to commit to really do form \
a legal Battleship fleet -- one five-cell carrier, one four-cell battleship, two three-cell \
cruisers, one two-cell destroyer, no overlaps, no diagonals, all in bounds -- AND their Pedersen \
hash equals this public commitment." That whole statement becomes one zk-SNARK. The proof is a \
few hundred bytes. I generate it in the browser before I ever send a transaction. The contract \
verifies it once and now the committed board is not just consistent -- it is provably LEGAL. \
The all-empty-board attack you just saw? It literally cannot produce a valid proof, because \
the circuit's fleet-histogram constraint rejects it. The proof IS the validity certificate.

[Pranav]
And from there shot responses are cheap -- we could keep using Merkle for them, or a tiny \
per-shot circuit. The expensive legality check only happens once.
"""
    set_notes(slide, notes)
    return slide


# ===========================================================================
# Slide 7 — Noir + UltraPlonk
# ===========================================================================

def slide_7_noir(prs):
    slide = blank_slide(prs)
    chrome(slide, 7, eyebrow=f"06  {MIDOT}  Toolchain",
           title="Noir + UltraPlonk.")

    # Left: why noir bullets
    left_l = MARGIN
    left_w = Inches(6.4)
    col_top = Inches(2.55)

    add_text(slide, left_l, col_top, left_w, Inches(0.4),
             "WHY NOIR",
             font=BODY_FONT, size=11, bold=True, color=CYAN, tracking=250)

    why = [
        (f"{ARROW}", "Rust-like DX, readable circuits."),
        (f"{ARROW}", "In-browser proving via @aztec/bb.js WASM."),
        (f"{ARROW}", "UltraPlonk backend: no trusted setup."),
        (f"{ARROW}", "Auto-generates a Solidity verifier contract."),
        (f"{ARROW}", "~250k gas per on-chain verify."),
    ]
    _draw_bullets(slide, why, left_l, col_top + Inches(0.5), left_w,
                  mark_color=ORANGE, text_color=TEXT)

    # Right: circuit stat card
    right_l = left_l + left_w + Inches(0.4)
    right_w = SLIDE_W - right_l - MARGIN

    add_round_rect(slide, right_l, col_top, right_w, Inches(4.2),
                   fill_color=SURFACE, line_color=HAIRLINE,
                   line_width=Pt(1), corner=0.08)

    add_text(slide, right_l + Inches(0.4), col_top + Inches(0.3),
             right_w - Inches(0.8), Inches(0.4),
             "board_validity.nr",
             font=CODE_FONT, size=12, bold=True, color=CYAN)
    add_text(slide, right_l + Inches(0.4), col_top + Inches(0.7),
             right_w - Inches(0.8), Inches(0.6),
             "3,482",
             font=BODY_FONT, size=42, bold=True, color=ORANGE)
    add_text(slide, right_l + Inches(0.4), col_top + Inches(1.45),
             right_w - Inches(0.8), Inches(0.4),
             "constraints",
             font=BODY_FONT, size=13, bold=False, color=MUTED)

    stats = [
        ("leaves",       "100 cells"),
        ("fleet",        "1x5  1x4  2x3  1x2"),
        ("hash",         "pedersen(cells, salt)"),
        ("backend",      "UltraPlonk / Honk"),
        ("verifier gas", "~250,000"),
    ]
    sy = col_top + Inches(2.0)
    for i, (k, v) in enumerate(stats):
        y = sy + Inches(0.35) * i
        add_text(slide, right_l + Inches(0.4), y,
                 Inches(1.8), Inches(0.3),
                 k,
                 font=CODE_FONT, size=11, bold=False, color=MUTED)
        add_text(slide, right_l + Inches(2.1), y,
                 right_w - Inches(2.5), Inches(0.3),
                 v,
                 font=CODE_FONT, size=11, bold=True, color=TEXT)

    notes = """\
[Vikram]
Quick tool choice. Noir is Aztec's Rust-flavoured DSL for writing zk circuits. We picked it \
over Circom for three reasons. One, the developer experience is a lot nicer -- it reads like \
normal code. Two, its UltraPlonk backend does not need a trusted setup ceremony, which means I \
do not have to spend five minutes of this talk explaining powers of tau. Three, it ships a \
browser prover as a WASM module so the whole proof lifecycle -- placement, prove, verify -- \
stays on one laptop. The circuit is called board_validity.nr, 3,482 constraints, and it's the \
exact thing I described on the previous slide: shapes, counts, no-diagonal run checks, and a \
Pedersen commitment binding the private board to the public hash.

[Pranav]
And Noir spits out a Solidity verifier contract we drop straight into the repo.
"""
    set_notes(slide, notes)
    return slide


# ===========================================================================
# Slide 8 — System Architecture
# ===========================================================================

def slide_8_architecture(prs):
    slide = blank_slide(prs)
    chrome(slide, 8, eyebrow=f"07  {MIDOT}  Stack",
           title="System architecture.")

    stations = [
        ("BRD", "Board",           "private[100]"),
        ("PSD", "Pedersen",        f"hash {ARROW} commit"),
        ("NOR", "Noir circuit",    "board_validity"),
        ("UPL", "UltraPlonk",      "proof blob"),
        ("VER", "Verifier",        "HonkVerifier.sol"),
        ("GME", "BattleshipGame",  "settle on-chain"),
    ]

    usable_w = SLIDE_W - MARGIN * 2
    tile_w   = Inches(1.85)
    tile_h   = Inches(1.4)
    n = len(stations)
    gap_w    = (usable_w - tile_w * n) / (n - 1)
    row_y    = Inches(2.75)

    centers = []
    for i, (glyph, title, sub) in enumerate(stations):
        x = MARGIN + (tile_w + gap_w) * i
        add_flow_station(slide, x, row_y, tile_w, tile_h, glyph, title, sub)
        centers.append((x, x + tile_w, row_y + tile_h / 2))

    for i in range(n - 1):
        _, x_end_prev, cy = centers[i]
        x_start_next, _, _ = centers[i + 1]
        y = int(cy)
        add_arrow(
            slide,
            x_end_prev + Emu(20000), y,
            x_start_next - Emu(20000), y,
            color=ORANGE,
        )

    cap_top = row_y + tile_h + Inches(0.45)
    add_text(
        slide, MARGIN, cap_top, usable_w, Inches(0.35),
        f"BROWSER  {ARROW}  WASM PROVER  {ARROW}  NOIR  {ARROW}  VERIFIER  {ARROW}  CHAIN",
        font=BODY_FONT, size=10, bold=True, color=MUTED, tracking=250,
        align=PP_ALIGN.CENTER,
    )

    bul_top = Inches(5.55)
    bullets = [
        (f"{MIDOT}", "Browser proves, chain verifies", CYAN),
        (f"{MIDOT}", "Burner wallets, zero popups",    ORANGE),
        (f"{MIDOT}", "Local Anvil for instant demo",   GOOD),
    ]
    col_w_b = usable_w / 3
    for i, (dot, text, color) in enumerate(bullets):
        x = MARGIN + col_w_b * i
        add_text(slide, x, bul_top, Inches(0.3), Inches(0.35),
                 dot,
                 font=BODY_FONT, size=18, bold=True, color=color)
        add_text(slide, x + Inches(0.3), bul_top - Inches(0.02),
                 col_w_b - Inches(0.3), Inches(0.4),
                 text,
                 font=BODY_FONT, size=14, bold=False, color=TEXT)

    notes = """\
[Pranav]
Data flow end to end. I open the React app. I drag my fleet onto the grid. I click Ready. The \
frontend runs bb.js, which is Aztec's WASM-compiled Barretenberg prover, against the Noir \
board-validity circuit right there in the browser. That spits out a proof blob. I send it as \
calldata to BattleshipGame.sol on a local Anvil chain. The contract calls the Noir-generated \
HonkVerifier, which accepts or rejects the proof on-chain. Every shot response goes through the \
same pipeline with the shot-response circuit. Nothing trusts the frontend -- every claim is \
verified by the contract. Two burner wallets derived from Anvil's deterministic test keys mean \
no MetaMask popups, no wallet switching, no friction. It is genuinely a one-laptop hot-seat \
demo and that is what you're about to see.
"""
    set_notes(slide, notes)
    return slide


# ===========================================================================
# Slide 9 — Live Demo
# ===========================================================================

def slide_9_demo(prs):
    slide = blank_slide(prs)

    add_text(
        slide, MARGIN, Inches(0.9),
        Inches(12.0), Inches(0.35),
        f"08  {MIDOT}  SWITCH TO THE BROWSER",
        font=BODY_FONT, size=11, bold=True, color=CYAN, tracking=350,
    )

    # "LIVE DEMO ->" — plain ASCII + arrow, no play triangle
    add_text(
        slide, Inches(0.0), Inches(2.35),
        Inches(13.333), Inches(2.3),
        f"LIVE DEMO  {ARROW}",
        font=BODY_FONT, size=96, bold=True, color=ORANGE,
        align=PP_ALIGN.CENTER,
    )

    add_text(
        slide, Inches(0.0), Inches(4.75),
        Inches(13.333), Inches(0.55),
        "http://localhost:5173",
        font=CODE_FONT, size=20, bold=False, color=TEXT,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        slide, Inches(0.0), Inches(5.25),
        Inches(13.333), Inches(0.45),
        f"two burner wallets  {MIDOT}  one laptop  {MIDOT}  zero trust",
        font=BODY_FONT, size=14, bold=False, color=MUTED,
        align=PP_ALIGN.CENTER, tracking=150,
    )

    add_wordmark(slide)
    add_progress_bar(slide, 9)

    notes = """\
[Pranav]
Alright, I'm driving. I'm opening the app -- you see both players' grids side by side. I'm \
clicking the carrier in the palette, rotating with R, dropping it on the top row. I'll place \
the rest fast. Hit Ready. Watch the spinner: "proving board legality" -- that is the Noir \
circuit running right now, 3,482 constraints in the browser. Proof verified on-chain. Player \
two does the same. Now we fire -- I click (3, 5) on the enemy grid. Pending pulse. Boom, hit \
confirmed on-chain, shot-response proof verified in about a second. Keep an eye on the Crypto \
Log panel on the right -- every on-chain event is narrated as it happens.

[Vikram]
Watch the sunk animation when the last cell of a ship goes down. This is a full cryptographic \
game loop running live with zero trust in the opponent.
"""
    set_notes(slide, notes)
    return slide


# ===========================================================================
# Slide 10 — What Was Proven
# ===========================================================================

def slide_10_proven(prs):
    slide = blank_slide(prs)
    chrome(slide, 10, eyebrow=f"09  {MIDOT}  Recap",
           title="What was proven on-chain.")

    items = [
        ("1", "Board legality",       "Before a single shot is fired."),
        ("2", "Shot-response honesty", "Hit or miss, without revealing the board."),
        ("3", "No trusted setup",      "UltraPlonk, no ceremony, no toxic waste."),
        ("4", "Cheap verification",    "About 250k gas per proof."),
    ]
    row_top = Inches(2.55)
    row_h   = Inches(0.95)
    circle  = Inches(0.7)

    for i, (num, title, body) in enumerate(items):
        y = row_top + row_h * i
        add_oval(slide, MARGIN, y + Inches(0.02),
                 circle, circle,
                 fill_color=ORANGE)
        add_text(slide, MARGIN, y + Inches(0.08),
                 circle, Inches(0.6),
                 num,
                 font=BODY_FONT, size=24, bold=True, color=BG,
                 align=PP_ALIGN.CENTER)
        add_text(slide, MARGIN + circle + Inches(0.3), y + Inches(0.02),
                 Inches(10), Inches(0.5),
                 title,
                 font=BODY_FONT, size=22, bold=True, color=TEXT)
        add_text(slide, MARGIN + circle + Inches(0.3), y + Inches(0.45),
                 Inches(10.5), Inches(0.4),
                 body,
                 font=BODY_FONT, size=14, bold=False, color=MUTED)

    notes = """\
[Vikram]
Four things the cryptography actually guaranteed during that demo. One: before the first shot, \
each player produced a zk-SNARK proving their board contained exactly the standard fleet -- \
seventeen cells, correct ship shapes, no overlaps, no diagonals. The contract either accepted \
that proof or the game did not start. Two: for every shot, the responding player produced a \
second proof binding the answer back to the original committed board, so they could never \
flip a hit to a miss. Three: none of this needed a trusted setup ceremony, because we used \
UltraPlonk. Four: on-chain verification cost about 250,000 gas per proof -- cheap enough to \
live inside a real game loop. No trust, no reveal, no ceremony.
"""
    set_notes(slide, notes)
    return slide


# ===========================================================================
# Slide 11 — What's Next / Q&A
# ===========================================================================

def slide_11_next(prs):
    slide = blank_slide(prs)
    chrome(slide, 11, eyebrow=f"10  {MIDOT}  Roadmap",
           title="What's next.")

    left_l = MARGIN
    col_top = Inches(2.55)
    left_w = Inches(6.6)

    steps = [
        ("01", "Mainnet deployment path",    "Solidity verifier is deploy-ready."),
        ("02", "Shot-response optimization", "Target sub-second proving."),
        ("03", "Multi-game lobby",           "Matchmaking plus state routing."),
        ("04", "Tournament mode",            "Trustless brackets, auditable replays."),
    ]
    for i, (num, title, body) in enumerate(steps):
        y = col_top + Inches(0.95) * i
        add_text(slide, left_l, y, Inches(0.6), Inches(0.4),
                 num,
                 font=CODE_FONT, size=13, bold=True, color=ORANGE, tracking=150)
        add_text(slide, left_l + Inches(0.7), y - Inches(0.03),
                 left_w - Inches(0.7), Inches(0.45),
                 title,
                 font=BODY_FONT, size=18, bold=True, color=TEXT)
        add_text(slide, left_l + Inches(0.7), y + Inches(0.38),
                 left_w - Inches(0.7), Inches(0.4),
                 body,
                 font=BODY_FONT, size=12, bold=False, color=MUTED)

    div_x = Inches(7.55)
    add_rect(slide, div_x, Inches(2.55), Pt(1), Inches(3.9),
             fill_color=HAIRLINE)

    right_l = Inches(7.95)
    add_text(slide, right_l, Inches(2.6),
             Inches(5.0), Inches(0.4),
             "Q & A",
             font=BODY_FONT, size=11, bold=True, color=CYAN, tracking=350)
    add_text(slide, right_l, Inches(3.1),
             Inches(5.2), Inches(2.0),
             "Questions?",
             font=BODY_FONT, size=60, bold=True, color=ORANGE)
    add_text(slide, right_l, Inches(4.85),
             Inches(5.2), Inches(1.4),
             "Circuits, contracts, frontend -- everything is open source.",
             font=BODY_FONT, size=14, bold=False, color=MUTED)
    add_text(slide, right_l, Inches(5.7),
             Inches(5.2), Inches(0.4),
             "github.com/ battleship-zk",
             font=CODE_FONT, size=12, bold=False, color=TEXT)

    notes = """\
[Pranav]
Where does this go from here? Mainnet deployment is basically a deploy script away -- the \
verifier is already Solidity. Shot-response proving time can come down with some circuit \
pruning. Longer term we want a multi-game lobby and a tournament mode where every move in \
every bracket is cryptographically auditable. That's the pitch. Thanks for your time.

[Vikram]
We're happy to take questions on the circuits, the contract, Noir, bb.js, the frontend -- \
anything you saw in the demo. The full stack is open source.
"""
    set_notes(slide, notes)
    return slide


# ===========================================================================
# Main
# ===========================================================================

def main():
    out_path = Path(__file__).parent / "battleship-zk-demo.pptx"

    prs = new_prs()
    slide_1_title(prs)
    slide_2_problem(prs)
    slide_3_merkle_101(prs)
    slide_4_merkle_battleship(prs)
    slide_5_merkle_hole(prs)
    slide_6_zksnarks(prs)
    slide_7_noir(prs)
    slide_8_architecture(prs)
    slide_9_demo(prs)
    slide_10_proven(prs)
    slide_11_next(prs)

    prs.save(str(out_path))
    print(f"Saved: {out_path}  ({out_path.stat().st_size:,} bytes)")
    print("Fonts referenced by name: Inter (body), JetBrains Mono (code).")
    print(f"  Fallbacks: Helvetica Neue {ARROW} Calibri (body), Menlo {ARROW} Consolas (code).")


if __name__ == "__main__":
    main()
